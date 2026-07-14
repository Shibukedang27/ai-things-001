from __future__ import annotations

import re
from io import BytesIO
from urllib.parse import urlparse

import httpx
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from knowledge_engine.exceptions import DocumentParsingError
from knowledge_engine.types import SourceType
from knowledge_engine.utils import normalize_space


class TextExtractor:
    def extract(self, content: bytes | str) -> str:
        if isinstance(content, str):
            return content
        return content.decode("utf-8", errors="ignore")


class PdfExtractor:
    def extract(self, content: bytes | str) -> str:
        if isinstance(content, str):
            content = content.encode("utf-8", errors="ignore")
        try:
            reader = PdfReader(BytesIO(content))
            pages = [page.extract_text() or "" for page in reader.pages]
        except Exception as exc:
            raise DocumentParsingError("PDF content could not be parsed.") from exc
        return "\n\n".join(page.strip() for page in pages if page.strip())


class DocxExtractor:
    def extract(self, content: bytes | str) -> str:
        if isinstance(content, str):
            content = content.encode("utf-8", errors="ignore")
        try:
            document = DocxDocument(BytesIO(content))
        except Exception as exc:
            raise DocumentParsingError("DOCX content could not be parsed.") from exc
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        table_text: list[str] = []
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_text.append(" | ".join(cells))
        return "\n\n".join([*paragraphs, *table_text])


class PowerPointExtractor:
    def extract(self, content: bytes | str) -> str:
        if isinstance(content, str):
            content = content.encode("utf-8", errors="ignore")
        try:
            presentation = Presentation(BytesIO(content))
        except Exception as exc:
            raise DocumentParsingError("PowerPoint content could not be parsed.") from exc

        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            fragments: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    fragments.append(shape.text.strip())
            if fragments:
                slides.append(f"Slide {index}\n" + "\n".join(fragments))
        return "\n\n".join(slides)


class MarkdownExtractor:
    CODE_FENCE_PATTERN = re.compile(r"```[\s\S]*?```")

    def extract(self, content: bytes | str) -> str:
        text = TextExtractor().extract(content)
        text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
        text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
        text = re.sub(r"^\s{0,3}#{1,6}\s*", "", text, flags=re.MULTILINE)
        return text


class WebsiteExtractor:
    def extract(self, url: str) -> tuple[str, dict[str, object]]:
        try:
            response = httpx.get(url, follow_redirects=True, timeout=15.0)
            response.raise_for_status()
        except httpx.HTTPError as exc:
            raise DocumentParsingError("Website URL could not be fetched.") from exc

        soup = BeautifulSoup(response.text, "html.parser")
        for tag in soup(["script", "style", "noscript", "svg", "canvas"]):
            tag.decompose()

        title = self._title(soup, url)
        author = self._meta_content(soup, "author")
        description = self._meta_content(soup, "description")
        content_nodes = soup.find_all(["h1", "h2", "h3", "p", "li", "blockquote"])
        body_parts = [node.get_text(" ", strip=True) for node in content_nodes]
        text = "\n".join(part for part in body_parts if part)
        metadata = {
            "url": url,
            "domain": urlparse(url).netloc,
            "title": title,
            "author": author,
            "description": description,
            "status_code": response.status_code,
        }
        return normalize_space(text), metadata

    def _title(self, soup: BeautifulSoup, url: str) -> str:
        if soup.title and soup.title.text.strip():
            return soup.title.text.strip()
        return urlparse(url).netloc or url

    def _meta_content(self, soup: BeautifulSoup, name: str) -> str | None:
        tag = soup.find("meta", attrs={"name": name}) or soup.find("meta", attrs={"property": f"og:{name}"})
        if tag and tag.get("content"):
            return str(tag["content"]).strip()
        return None


class GitHubRepositoryExtractor:
    README_CANDIDATES = ("README.md", "readme.md", "README.rst", "README.txt")

    def extract(self, url: str) -> tuple[str, dict[str, object]]:
        owner, repo = self._parse_repo(url)
        metadata = {
            "url": url,
            "owner": owner,
            "repository": repo,
            "source": "github",
        }
        for branch in ("main", "master"):
            for candidate in self.README_CANDIDATES:
                raw_url = f"https://raw.githubusercontent.com/{owner}/{repo}/{branch}/{candidate}"
                try:
                    response = httpx.get(raw_url, follow_redirects=True, timeout=15.0)
                except httpx.HTTPError:
                    continue
                if response.status_code == 200 and response.text.strip():
                    metadata.update({"readme_url": raw_url, "branch": branch, "readme_file": candidate})
                    return response.text, metadata
        raise DocumentParsingError("GitHub repository README could not be fetched from main or master.")

    def _parse_repo(self, url: str) -> tuple[str, str]:
        parsed = urlparse(url)
        parts = [part for part in parsed.path.strip("/").split("/") if part]
        if parsed.netloc.lower() != "github.com" or len(parts) < 2:
            raise DocumentParsingError("GitHub repository URL must look like https://github.com/{owner}/{repo}.")
        return parts[0], parts[1].removesuffix(".git")


EXTRACTORS_BY_SOURCE_TYPE = {
    SourceType.PDF: PdfExtractor(),
    SourceType.DOCX: DocxExtractor(),
    SourceType.TXT: TextExtractor(),
    SourceType.MARKDOWN: MarkdownExtractor(),
    SourceType.POWERPOINT: PowerPointExtractor(),
    SourceType.RESEARCH_PAPER: PdfExtractor(),
    SourceType.YOUTUBE_TRANSCRIPT: TextExtractor(),
    SourceType.PLAIN_NOTES: TextExtractor(),
}
